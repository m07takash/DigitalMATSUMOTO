import os
import logging
import pytz
import re
import numpy as np
import json
import yaml
import shutil
import mimetypes
import math
import bcrypt
import pdfplumber
from dateutil import parser
from datetime import datetime
from dotenv import load_dotenv

import unicodedata
from pathlib import Path

import MeCab
from sklearn.feature_extraction.text import TfidfVectorizer
from wordcloud import WordCloud

import base64
import tiktoken
import openai
from openai import OpenAI, AzureOpenAI
from google import genai
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Emu

# Load system.env and set environment variables
if os.path.exists("system.env"):
    load_dotenv("system.env")
timezone = os.getenv("TIMEZONE")
openai_api_key = os.getenv("OPENAI_API_KEY")
gemini_api_key = os.getenv("GEMINI_API_KEY")
embedding_model = os.getenv("EMBEDDING_MODEL")
# Embedding provider switch ("openai"=OpenAI, "azure"=Azure OpenAI)
embed_provider = (os.getenv("EMBED_PROVIDER") or "openai").lower()
azure_openai_embed_model = os.getenv("AZURE_OPENAI_EMBED_MODEL")  # Azure deployment name

# Logging initialization (call once at application startup)
def setup_logging(level=logging.INFO):
    logging.basicConfig(
        level=level,
        format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
        datefmt="%Y-%m-%d %H:%M:%S"
    )

# Convert a timestamp string to a time
def safe_parse_timestamp(timestamp_str):
    jst = pytz.timezone(timezone)
    try:
        return jst.localize(datetime.strptime(timestamp_str, "%Y/%m/%d %H:%M:%S")).isoformat()
    except ValueError:
        return datetime.now(jst).isoformat()

# Convert a date-style string
def parse_date(d):
    try:
        return datetime.fromisoformat(d)
    except Exception:
        return datetime.min

# Safely parse a LOG_TEMPLATE-formatted string into a dict
def parse_log_template(ref_str: str) -> dict:
    """Convert a LOG_TEMPLATE-formatted string ('key': value, ...) to a dict.
    Robust to single quotes, double quotes, and special characters inside the text."""
    import re
    s = str(ref_str).replace("\n", "").replace("$", "＄")
    try:
        return ast.literal_eval("{" + s + "}")
    except Exception:
        pass
    # Fallback: split by known key names and parse
    result = {}
    # Find every 'key': pattern position
    key_pattern = re.compile(r"'(\w+)'\s*:\s*")
    matches = list(key_pattern.finditer(s))
    for i, m in enumerate(matches):
        key = m.group(1)
        val_start = m.end()
        val_end = matches[i + 1].start() if i + 1 < len(matches) else len(s)
        raw_val = s[val_start:val_end].strip().rstrip(",").strip()
        # Strip quotes
        if (raw_val.startswith("'") and raw_val.endswith("'")) or \
           (raw_val.startswith('"') and raw_val.endswith('"')):
            raw_val = raw_val[1:-1]
        # Try to coerce to a number
        try:
            result[key] = int(raw_val)
        except (ValueError, TypeError):
            try:
                result[key] = float(raw_val)
            except (ValueError, TypeError):
                result[key] = raw_val
    return result

# Sanitize text (remove JSON/XML-forbidden and control characters)
def sanitize_text(text: str) -> str:
    """Strip control characters that would break JSON/XML in LLM response text.
    Surrogate pairs (e.g. emoji) are preserved."""
    if not isinstance(text, str):
        return str(text) if text is not None else ""
    # Strip NUL and control characters (keep tab / newline / CR)
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)

# Sanitize file names
def sanitize_filename(name: str, replacement: str = "_", max_length: int = 255, keep_unicode: bool = True) -> str:
    """
    Return a file name with dangerous/inconvenient characters removed or replaced.
    - Strip path separators /, \\ and control characters
    - Strip Windows-forbidden characters <>:"/\\|?*
    - Also remove ", ', $, %, / shown in examples
    - Remove trailing dots/spaces (Windows safety)
    - When the result is empty, return 'untitled'
    """
    if name is None:
        return "untitled"
    
    # Windows reserved names (cannot/should not save when the stem matches)
    _WINDOWS_RESERVED = {
        "CON", "PRN", "AUX", "NUL",
        *(f"COM{i}" for i in range(1, 10)),
        *(f"LPT{i}" for i in range(1, 10)),
    }

    # Surrounding whitespace (trailing whitespace causes issues on Windows)
    name = name.strip()

    # Unicode normalization (align look-alike characters)
    name = unicodedata.normalize("NFKC", name)

    # Character handling: pass keep_unicode=False to flatten to ASCII
    if not keep_unicode:
        name = name.encode("ascii", "ignore").decode("ascii")

    # Strip control characters
    name = "".join(ch for ch in name if unicodedata.category(ch)[0] != "C")

    # Always strip path separators (mitigates directory traversal)
    name = name.replace("/", replacement).replace("\\", replacement)

    # Replace/remove Windows-forbidden characters: <>:"/\|?*
    name = re.sub(r'[<>:"/\\|?*]', replacement, name)

    # Strip a broader set of awkward symbols (as in the docstring examples)
    # Adjust this set as needed
    name = re.sub(r"""["'$%]""", "", name)

    # Collapse consecutive replacement characters
    if replacement:
        rep_esc = re.escape(replacement)
        name = re.sub(rf"{rep_esc}+", replacement, name)

    # Trim leading/trailing replacements and dots/spaces
    name = name.strip(" ._")

    # Fallback if the result is empty
    if not name:
        name = "untitled"

    # Windows reserved-name guard: if the stem is reserved, prefix with _
    stem, dot, suffix = name.partition(".")
    if stem.upper() in _WINDOWS_RESERVED:
        name = f"_{name}"

    # If too long (preserve the extension when possible)
    if len(name) > max_length:
        p = Path(name)
        ext = p.suffix
        base = p.stem
        cut = max_length - len(ext)
        name = base[:max(1, cut)] + ext

    name = name.rstrip(" .")
    return name

# Password hashing and verification
def hash_password(plain: str) -> str:
    # bcrypt operates on bytes
    salt = bcrypt.gensalt(rounds=12)
    hashed = bcrypt.hashpw(plain.encode("utf-8"), salt)
    return hashed.decode("utf-8")

# Verify a password
def verify_password(plain: str, stored: str) -> bool:
    if stored.startswith("$2a$") or stored.startswith("$2b$") or stored.startswith("$2y$"):
        return bcrypt.checkpw(plain.encode("utf-8"), stored.encode("utf-8"))
    return plain == stored

# Convert a local file into the same object type as Streamlit's UploadedFile
class ConvertedLocalFile:
    def __init__(self, file_path):
        self.file_path = file_path
        self.name = os.path.basename(file_path)
        self.size = os.path.getsize(file_path)
        self.type, _ = mimetypes.guess_type(file_path)

    def read(self):
        with open(self.file_path, 'rb') as f:
            return f.read()

# Helper to emit a list from a pattern
def extract_list_pattern(text, pattern=r"(\[\s*{.*?}\s*\])"):
    match = re.search(pattern, text, re.DOTALL)
    if match:
        json_str = match.group(1)
        try:
            return json.loads(json_str)
        except json.JSONDecodeError:
            return []
    return []

# Trim a list of date ranges
def merge_periods(periods):
    # The upstream extract_date LLM occasionally returns items that lack one
    # or both date fields (e.g. `{"end": "2026/06/20"}` or `{"label": "今"}`).
    # Silently skip those rather than KeyError-ing the entire turn — Meta
    # Search is supplemental, not load-bearing.
    converted = []
    for p in periods:
        if not isinstance(p, dict):
            continue
        _s = p.get("start")
        _e = p.get("end")
        if not _s or not _e:
            continue
        try:
            converted.append({
                "start": datetime.strptime(str(_s), "%Y/%m/%d"),
                "end":   datetime.strptime(str(_e), "%Y/%m/%d"),
            })
        except (ValueError, TypeError):
            continue

    # Sort by start
    converted.sort(key=lambda x: x["start"])
    merged = []
    for p in converted:
        if not merged:
            merged.append(p)
            continue
        last = merged[-1]

        # Merge when overlapping or adjacent
        if p["start"] <= last["end"]:
            last["end"] = max(last["end"], p["end"])
        else:
            merged.append(p)

    # datetime
    result = [
        {
            "start": m["start"].strftime("%Y/%m/%d"),
            "end": m["end"].strftime("%Y/%m/%d"),
        }
        for m in merged
    ]

    return result

# Recursively overwrite a dictionary
def update_dict(target, source):
    for key, value in source.items():
        if isinstance(value, dict) and key in target and isinstance(target[key], dict):
            # Update recursively when a nested dict exists
            update_dict(target[key], value)
        else:
            target[key] = value

# Fetch multiple files from a folder
def get_files(folder_path="", identifier="", file_sort="Y"):
    files = []
    if os.path.exists(folder_path):
        all_files = os.listdir(folder_path)
        files = [f for f in all_files if f.endswith(identifier)]
        if file_sort == "Y":
            files.sort()
    return files

# Move a file
def copy_file(from_file_path, to_file_path):
    shutil.move(from_file_path, to_file_path)

# Read a text file
def read_text_file(file_name, folder_path=""):
    file_path = folder_path + file_name
    if os.path.exists(file_path):
        with open(file_path, 'r') as file:
            data = file.read()
    else:
        data = ""
    return data

# Read a PDF file
def read_pdf_file(file_name, folder_path=""):
    file_path = folder_path + file_name
    pdf_dict = {}
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            # Extract each page's text and store per-page in a dict
            pdf_dict[f"Page_{i + 1}"] = page.extract_text()
    return pdf_dict

# Extract text + images from a PDF
def read_pdf_with_images(file_path, temp_folder):
    """Extract per-page text and embedded images from a PDF."""
    pdf_text = {}
    image_files = []
    os.makedirs(temp_folder, exist_ok=True)

    # Text extraction (pdfplumber)
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            # Also extract tables
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    text += "\n[Table]\n"
                    for row in table:
                        text += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
            pdf_text[f"Page_{i + 1}"] = text

    # Image extraction (PyMuPDF)
    doc = fitz.open(file_path)
    img_idx = 0
    for page_num in range(len(doc)):
        page = doc[page_num]
        images = page.get_images(full=True)
        for img in images:
            xref = img[0]
            base_image = doc.extract_image(xref)
            if base_image and base_image["image"]:
                ext = base_image.get("ext", "png")
                img_path = os.path.join(temp_folder, f"pdf_p{page_num+1}_img{img_idx}.{ext}")
                with open(img_path, "wb") as f:
                    f.write(base_image["image"])
                image_files.append(img_path)
                img_idx += 1
    doc.close()
    return pdf_text, image_files

# Extract text + images from DOCX
def read_docx_file(file_path, temp_folder):
    """Extract paragraph text, tables, and embedded images from a DOCX."""
    os.makedirs(temp_folder, exist_ok=True)
    doc = DocxDocument(file_path)
    text_parts = []
    image_files = []
    img_idx = 0

    # Paragraph text
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    # Tables
    for table in doc.tables:
        text_parts.append("[Table]")
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            text_parts.append(" | ".join(cells))

    # Embedded images
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            img_data = rel.target_part.blob
            ext = os.path.splitext(rel.target_ref)[1] or ".png"
            img_path = os.path.join(temp_folder, f"docx_img{img_idx}{ext}")
            with open(img_path, "wb") as f:
                f.write(img_data)
            image_files.append(img_path)
            img_idx += 1

    return "\n".join(text_parts), image_files

# Extract text from XLSX
def read_xlsx_file(file_path):
    """Extract cell data from every sheet in an XLSX in table form."""
    import openpyxl
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        text_parts.append(f"[Sheet: {sheet_name}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            if any(cells):
                text_parts.append(" | ".join(cells))
    return "\n".join(text_parts)

# Extract text + images from PPTX
def read_pptx_file(file_path, temp_folder):
    """Extract per-slide text, tables, notes, and embedded images from a PPTX."""
    os.makedirs(temp_folder, exist_ok=True)
    prs = Presentation(file_path)
    text_parts = []
    image_files = []
    img_idx = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = f"\n[Slide {slide_num}]"

        # Title
        if slide.shapes.title and slide.shapes.title.text:
            slide_text += f"\nTitle: {slide.shapes.title.text}"

        # Text / table / image
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_text += f"\n{para.text}"
            elif shape.has_table:
                slide_text += "\n[Table]"
                for row in shape.table.rows:
                    cells = [cell.text for cell in row.cells]
                    slide_text += "\n" + " | ".join(cells)
            if shape.shape_type == 13:  # Picture
                try:
                    img_blob = shape.image.blob
                    ext = shape.image.content_type.split("/")[-1] if shape.image.content_type else "png"
                    img_path = os.path.join(temp_folder, f"pptx_s{slide_num}_img{img_idx}.{ext}")
                    with open(img_path, "wb") as f:
                        f.write(img_blob)
                    image_files.append(img_path)
                    img_idx += 1
                    slide_text += f"\n[Image: pptx_s{slide_num}_img{img_idx-1}.{ext}]"
                except Exception:
                    pass

        # Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_text += f"\nNotes: {notes}"

        text_parts.append(slide_text)

    return "\n".join(text_parts), image_files

# Read a JSON file.
# Concurrent writers (e.g., parallel persona threads writing chat_memory.json) write atomically
# via temp + rename, but to be safe against pathological races on non-POSIX filesystems we
# retry once on JSONDecodeError after a short sleep.
def read_json_file(file_name, folder_path=""):
    import time as _time
    data = {}
    # Reject empty file names — otherwise `folder_path + ""` returns the
    # folder itself, `os.path.exists` reports True, and the subsequent
    # `open()` raises IsADirectoryError. Callers that pass an empty name
    # (rare, but happens when a session has no agent set) should get an
    # empty dict back, matching the missing-file semantics.
    if not file_name:
        return data
    file_path = folder_path + file_name
    if not os.path.exists(file_path) or os.path.isdir(file_path):
        return data
    for _attempt in range(2):
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
            return data
        except json.JSONDecodeError:
            if _attempt == 0:
                _time.sleep(0.05)
                continue
            raise


# Atomic JSON write: write to a temp sibling, fsync, then rename.
# On POSIX (Linux/Mac), rename swaps the inode atomically, so concurrent readers either
# see the previous file in full or the new file in full -- never a partial one.
def save_json_file(data, file_path, indent=2):
    import tempfile
    folder = os.path.dirname(file_path) or "."
    os.makedirs(folder, exist_ok=True)
    fd, tmp_path = tempfile.mkstemp(prefix=".tmp_", suffix=".json", dir=folder)
    try:
        with os.fdopen(fd, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=indent)
            f.flush()
            try:
                os.fsync(f.fileno())
            except OSError:
                pass
        os.replace(tmp_path, file_path)
    except Exception:
        # Clean up the temp file on any failure so we don't litter the folder.
        try:
            os.remove(tmp_path)
        except OSError:
            pass
        raise

# Read a YAML file
def read_yaml_file(file_name, folder_path=""):
    # See `read_json_file` for why the retry is here: although writers go through
    # temp + rename, very rare partial reads can still occur on edge-case filesystems.
    import time as _time
    data = {}
    # Same guard as `read_json_file`: reject empty file names so
    # `folder_path + ""` doesn't accidentally open a directory.
    if not file_name:
        return data
    file_path = folder_path + file_name
    if not os.path.exists(file_path) or os.path.isdir(file_path):
        return data
    for _attempt in range(2):
        try:
            with open(file_path, 'r', encoding="utf-8") as file:
                data = yaml.safe_load(file)
            break
        except yaml.YAMLError:
            if _attempt == 0:
                _time.sleep(0.05)
                continue
            raise
    if data is None:
        data = {}
    return data

# Save a YAML file
def save_text_file(data, file_path):
    with open(file_path, 'w', encoding="utf-8") as f:
        f.write(data)

# Atomic + lock-protected YAML write.
# This file path is read-modify-write (existing keys are preserved; new keys are merged in),
# so we MUST hold the per-path lock during the whole read-merge-write cycle, otherwise
# concurrent writers race and lose updates from one another.
# The actual write itself goes via temp + rename so concurrent readers never see a partial file.
_yaml_file_locks = {}
_yaml_file_locks_meta = __import__("threading").Lock()


def _get_yaml_file_lock(file_path):
    import threading as _th
    with _yaml_file_locks_meta:
        if file_path not in _yaml_file_locks:
            _yaml_file_locks[file_path] = _th.Lock()
        return _yaml_file_locks[file_path]


def save_yaml_file(data, file_name, folder_path=""):
    import tempfile
    file_path = folder_path + file_name
    folder = os.path.dirname(file_path) or "."
    os.makedirs(folder, exist_ok=True)
    with _get_yaml_file_lock(file_path):
        # Re-read under the lock so we don't drop concurrent updates.
        current_data = read_yaml_file(file_name, folder_path)
        current_data.update(data)
        fd, tmp_path = tempfile.mkstemp(prefix=".tmp_", suffix=".yaml", dir=folder)
        try:
            with os.fdopen(fd, 'w', encoding='utf-8') as f:
                yaml.dump(current_data, f, allow_unicode=True)
                f.flush()
                try:
                    os.fsync(f.fileno())
                except OSError:
                    pass
            os.replace(tmp_path, file_path)
        except Exception:
            try:
                os.remove(tmp_path)
            except OSError:
                pass
            raise

# Convert an MP3 file to text (switch between OpenAI / Azure via TRANSCRIBE_PROVIDER)
def mp3_to_text(file_name, folder_path=""):
    _provider = (os.getenv("TRANSCRIBE_PROVIDER") or "openai").lower()
    if _provider == "azure":
        client = AzureOpenAI(
            api_key=os.getenv("AZURE_OPENAI_API_KEY"),
            azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
            api_version=os.getenv("AZURE_OPENAI_API_VERSION", "2024-12-01-preview"),
        )
        model_name = os.getenv("AZURE_OPENAI_WHISPER_MODEL") or "whisper"
    else:
        client = OpenAI(api_key=openai_api_key)
        model_name = "whisper-1"
    file_path = folder_path + file_name
    if os.path.exists(file_path):
        with open(file_path, "rb") as audio_file:
            transcript = client.audio.transcriptions.create(
                model=model_name,
                file=audio_file,
                language="ja"
            )
        data = transcript.text
    else:
        data = ""
    return data

# Base64-encode an image file　
def encode_image_file(image_path):
    image_base64 = None
    with open(image_path, "rb") as image_file:
        image_base64 = base64.b64encode(image_file.read()).decode('utf-8')
    return image_base64

# Return the difference between two timestamps as a string
def get_time_diff(timestamp1, timestamp2, format_str="%Y-%m-%d %H:%M:%S.%f"):
    time1 = datetime.strptime(timestamp1, format_str)
    time2 = datetime.strptime(timestamp2, format_str)
    time_difference = time2 - time1
    return str(time_difference)

# Convert a date type
def convert_to_ymd(date_str, date_format='%Y-%m-%d'):
    try:
        if '/' in date_str:
            # Parse a date that contains a slash
            parsed_date = datetime.strptime(date_str, '%Y/%m/%d')
        else:
            # Auto-parse a date without slashes
            parsed_date = parser.parse(date_str)
        # Format using the requested layout
        formatted_date = parsed_date.strftime(date_format)
        return formatted_date
    except ValueError:
        return "Invalid date format"

# Singleton embedding client / tokenizer
_embed_client = None
_embed_enc = None


class _ApproxEncoder:
    """Approximation encoder used when tiktoken is disabled.
    Used only to (a) clip embedding input to the token cap and (b) display token counts,
    so under a worst-case CJK assumption we treat 1 character = 2 tokens, meaning the
    existing `max_tokens=8192` check effectively clips at 4096 characters.
    Tokens are stored as (char, "") pairs so encode/decode reproduces the original text.
    """
    def encode(self, text):
        text = text or ""
        out = []
        for c in text:
            out.append(c)
            out.append("")
        return out

    def decode(self, tokens):
        # Picking the even-indexed elements of each pair restores the original text
        return "".join(tokens[::2])


def _tiktoken_disabled():
    return (os.getenv("TIKTOKEN_DISABLE") or "").lower() in ("true", "1", "yes", "on")


def _load_tiktoken_encoding(name_or_model, by_model=False):
    """Common helper that loads a tiktoken encoding.
    On closed networks the HTTPS fetch from openaipublic.blob.core.windows.net fails,
    so cache files (cl100k_base / o200k_base, etc.) must be placed at TIKTOKEN_CACHE_DIR
    beforehand. When network failure is detected, raise a RuntimeError that includes
    the remediation steps.
    """
    try:
        if by_model:
            return tiktoken.encoding_for_model(name_or_model)
        return tiktoken.get_encoding(name_or_model)
    except KeyError:
        # Unknown model name (e.g. Azure deployment name) -> fall back to cl100k_base (text-embedding-3 compatible)
        if by_model:
            return _load_tiktoken_encoding("cl100k_base", by_model=False)
        raise
    except Exception as e:
        _cache_dir = os.getenv("TIKTOKEN_CACHE_DIR") or "(not set)"
        raise RuntimeError(
            f"Failed to load tiktoken encoding '{name_or_model}': {type(e).__name__}: {e}\n"
            f"  Cause: tiktoken fetches the BPE file from https://openaipublic.blob.core.windows.net/encodings/ on first use.\n"
            f"         Closed networks cannot reach that host, which is why this fails.\n"
            f"  Current TIKTOKEN_CACHE_DIR = {_cache_dir}\n"
            f"  Remedy: on a machine with outbound connectivity, run the command below to build the cache,\n"
            f"          copy the resulting files to any directory on this VM,\n"
            f"          and add TIKTOKEN_CACHE_DIR=<that path> to system.env and restart.\n"
            f"        TIKTOKEN_CACHE_DIR=/tmp/tiktoken_cache python3 -c \"import tiktoken;\\\n"
            f"          [tiktoken.get_encoding(n) for n in ('cl100k_base','o200k_base','p50k_base','r50k_base')]\""
        ) from e

def _get_embed_client():
    """Embedding client. Uses Azure OpenAI when EMBED_PROVIDER=azure."""
    global _embed_client, _embed_enc
    if _embed_client is None:
        if embed_provider == "azure":
            _embed_client = AzureOpenAI(
                api_key=os.getenv("AZURE_OPENAI_API_KEY"),
                azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
                api_version=os.getenv("AZURE_OPENAI_API_VERSION", "2024-12-01-preview"),
            )
        else:
            _embed_client = OpenAI(api_key=openai_api_key)
    if _embed_enc is None:
        # When TIKTOKEN_DISABLE=true, bypass tiktoken entirely and use a char-length approximation (closed-env friendly)
        if _tiktoken_disabled():
            _embed_enc = _ApproxEncoder()
        else:
            # tiktoken expects an OpenAI model name. Since the Azure deployment name is unknown,
            # look up by EMBEDDING_MODEL (the real model name); on failure fall back to cl100k_base (text-embedding-3 compatible)
            _embed_enc = _load_tiktoken_encoding(embedding_model or "text-embedding-3-large", by_model=True)
    return _embed_client, _embed_enc

def _embed_model_name():
    """The name to pass as the API's model argument. Deployment name on Azure; model name on OpenAI."""
    if embed_provider == "azure":
        return azure_openai_embed_model or embedding_model
    return embedding_model

# Convert text to an embedding (selects OpenAI / Azure OpenAI based on EMBED_PROVIDER)
def embed_text(text):
    client, enc = _get_embed_client()
    max_tokens = 8192
    tokens = enc.encode(text)
    if len(tokens) > max_tokens:
        tokens = tokens[:max_tokens]
    safe_text = enc.decode(tokens)
    response = client.embeddings.create(model=_embed_model_name(), input=safe_text)
    return response.data[0].embedding

# C-3: Vectorize multiple texts with a single API call (batched)
def embed_texts_batch(texts, batch_token_limit=250000):
    client, enc = _get_embed_client()
    max_tokens = 8192
    safe_texts = []
    for text in texts:
        tokens = enc.encode(text)
        if len(tokens) > max_tokens:
            tokens = tokens[:max_tokens]
        safe_texts.append(enc.decode(tokens))

    # Send batches split by the token cap
    all_embeddings = [None] * len(safe_texts)
    batch_start = 0
    while batch_start < len(safe_texts):
        batch_tokens = 0
        batch_end = batch_start
        for i in range(batch_start, len(safe_texts)):
            t = len(enc.encode(safe_texts[i]))
            if batch_tokens + t > batch_token_limit and i > batch_start:
                break
            batch_tokens += t
            batch_end = i + 1
        batch = safe_texts[batch_start:batch_end]
        response = client.embeddings.create(model=_embed_model_name(), input=batch)
        for item in response.data:
            all_embeddings[batch_start + item.index] = item.embedding
        batch_start = batch_end

    return all_embeddings

# Save an array of embedding vectors to a single .npy file
def save_vectext_to_npy(vec_text, file_path, dtype="float32"):
    arr = np.asarray(vec_text, dtype=dtype)
    np.save(file_path, arr)

# Load an array of embedding vectors from a .npy file
def read_vectext_to_npy(file_path, mmap=False):
    # Return an empty array if the file does not exist
    if not os.path.exists(file_path):
        return np.array([])

    # When mmap is True, load as a memory-mapped array
    if mmap:
        return np.load(file_path, mmap_mode="r")
    else:
        return np.load(file_path)

# Morphological analysis (Owakati)
def tokenize_Owakati(text, mode="Default", stop_words=[], grammer=('名詞', '動詞', '形容詞', '副詞')):
    mecab = MeCab.Tagger("-Owakati")
    wakati_text = mecab.parse(text)
    tokens = wakati_text.split()
    if mode == "All":
        # Drop stop words
        tokens = [word for word in tokens if word not in stop_words]
    else:
        mecab = MeCab.Tagger()  # Create another MeCab instance with default options to obtain morphological tags
        valid_tokens = []
        for token in tokens:
            node = mecab.parseToNode(token)  # Run morphological analysis on each token
            while node:
                pos = node.feature.split(",")[0]
                if pos.startswith(grammer) and token not in stop_words:
                    valid_tokens.append(token)
                    break  # Exit the loop to avoid double-checking the same token
                node = node.next
        tokens = valid_tokens
    return tokens

# Build a TF-IDF vectorizer from a set of texts
def fit_tfidf(texts, mode="Default", stop_words=[], grammer=('名詞', '動詞', '形容詞', '副詞')):
    vectorizer = TfidfVectorizer(tokenizer=lambda x: tokenize_Owakati(x, mode, stop_words, grammer), lowercase=False)
    vectorizer.fit(texts)
    return vectorizer

# Get the TF-IDF value for a keyword
def get_tfidf(text, vectorizer):
    tfidf_matrix = vectorizer.transform([text])
    feature_names = vectorizer.get_feature_names_out()
    tfidf_scores = tfidf_matrix.toarray()[0]
    tfidf_pairs = sorted(zip(feature_names, tfidf_scores), key=lambda x: x[1], reverse=True)
    tfidf_dict = dict(tfidf_pairs)
    #tfidf_pairs = [(word, score) for word, score in tfidf_dict.items()]
    return tfidf_pairs, tfidf_dict

# Get TF-IDF from a list
def get_tfidf_list(text, vectorizer, TopN):
    tfidf_pairs, tfidf_dict = get_tfidf(text, vectorizer)
    tfidf_topN = tfidf_pairs[:TopN]  # Take the top-N items
    tfidf_topN_list = [i[0] for i in tfidf_topN]
    tfidf_topN_str = '、'.join([f'{item[0]}：{round(item[1], 10)}' for item in tfidf_topN])
    return tfidf_dict, tfidf_topN, tfidf_topN_str

# Visualize as a word cloud
def get_wordcloud(title, dict, folder_path="user/common/analytics/insight/"):
    font_path = "/usr/share/fonts/truetype/fonts-japanese-gothic.ttf"
    wc = WordCloud(background_color="white", font_path=font_path, width=800, height=600, max_words=50, contour_color='steelblue')
    wc.generate_from_frequencies(dict)
    file_name = folder_path + f"WordCloud_{title}.png"
    wc.to_file(file_name)

# Token counting
def count_token(tokenizer, model, text):
    tokens = 0
    if tokenizer == "tiktoken":
        if _tiktoken_disabled():
            # Approximate as char_count * 2 (CJK worst-case). Display only; precision is not required.
            tokens = len(text or "") * 2
        else:
            encoding = _load_tiktoken_encoding("cl100k_base")
            tokens = len(encoding.encode(text))
    elif tokenizer == "gemini":
        gemini_client = genai.Client(api_key=gemini_api_key)
        response_tokens = gemini_client.models.count_tokens(model=model, contents=text)
        tokens = response_tokens.total_tokens
    else:
        tokens = len(text)
    return tokens

# Similarity computation (cosine distance)
def calculate_cosine_distance(vec1, vec2):
    # Compute cosine similarity
    dot_product = sum(p*q for p, q in zip(vec1, vec2))
    magnitude_vec1 = math.sqrt(sum(p**2 for p in vec1))
    magnitude_vec2 = math.sqrt(sum(q**2 for q in vec2))
    if magnitude_vec1 == 0 or magnitude_vec2 == 0:
        # When either vector has zero magnitude, similarity is undefined
        return 0
    cosine_similarity = dot_product / (magnitude_vec1 * magnitude_vec2)
    # Compute cosine distance (1 - cosine similarity)
    cosine_distance = 1-cosine_similarity
    return cosine_distance

# Similarity computation (Minkowski distance)
def calculate_minkowski_distance(vec1, vec2, p):
    if p == 0:
        raise ValueError("p cannot be zero")
    elif p == float('inf'):
        return max(abs(a - b) for a, b in zip(vec1, vec2))
    else:
        return sum(abs(a - b) ** p for a, b in zip(vec1, vec2)) ** (1 / p)

# Penalty (linear)
def linear_penalty(difference, alpha=0.001):
    return alpha * difference

# Penalty (exponential)
def exponential_penalty(difference, alpha=0.001):
    return 1 - math.exp(-alpha * difference)

# Penalty (logistic)
def logistic_penalty(difference, alpha=0.001):
    return 1 / (1 + math.exp(-alpha * difference))

# Penalty (step)
def step_penalty(difference, alpha=0.001, beta=0.01, days=10):
    q = difference // days  # Quotient when dividing days by `days`
    r = difference % days  # Remainder when dividing days by `days`
    return q*beta

# Penalty (alpha increases per step) -- pick a base function separately (linear_penalty / exponential_penalty / logistic_penalty)
def step_gain(difference, alpha=0.001, days=10, func="linear_penalty"):
    quotient = difference // days  # Divide days by `days`
    alpha = quotient*alpha
    return globals()[func](difference, alpha)

# Similarity (distance + date penalty)
def calculate_similarity_vec(vec1, vec2, logic="Cosine"):
    if logic == "Cosine":
        distance = calculate_cosine_distance(vec1, vec2)
    elif logic == "Euclidean":
        distance = calculate_minkowski_distance(vec1, vec2, 2)
    elif logic == "Manhattan":
        distance = calculate_minkowski_distance(vec1, vec2, 1)
    elif logic == "Chebyshev":
        distance = calculate_minkowski_distance(vec1, vec2, float('inf'))
    else:  # default to cosine distance
        distance = calculate_cosine_distance(vec1, vec2)
    return distance